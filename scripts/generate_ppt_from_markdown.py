#!/usr/bin/env python3
"""Generate a .pptx from the Final Presentation markdown file.

Usage:
  python scripts/generate_ppt_from_markdown.py \
    --input "Project documentation/Final Presentation - Pantry Planner.md" \
    --output "Project documentation/Final Presentation - Pantry Planner.pptx"
"""

from __future__ import annotations

import argparse
import re
from pathlib import Path


SLIDE_HEADER_RE = re.compile(r"^##\s+Slide\s+\d+\s+—\s+(.*)$")


def parse_markdown_sections(markdown_text: str):
    sections = []
    current_title = None
    current_lines = []

    for raw_line in markdown_text.splitlines():
        line = raw_line.rstrip()
        header_match = SLIDE_HEADER_RE.match(line)

        if header_match:
            if current_title is not None:
                sections.append((current_title, current_lines))
            current_title = header_match.group(1).strip()
            current_lines = []
            continue

        if current_title is None:
            continue

        if line.startswith("**Speaker:**"):
            continue
        if line.strip() == "---":
            continue

        # Light cleanup of markdown marks.
        cleaned = line.replace("**", "").replace("`", "")
        if cleaned.startswith("> "):
            cleaned = cleaned[2:]
        current_lines.append(cleaned)

    if current_title is not None:
        sections.append((current_title, current_lines))

    return sections


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_content_slide(prs: Presentation, title: str, lines: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()

    first_written = False
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue

        bullet_level = 0
        if stripped.startswith("- "):
            text = stripped[2:].strip()
            bullet_level = 0
        elif stripped.startswith(tuple(f"{n}. " for n in range(1, 10))):
            text = stripped
            bullet_level = 0
        elif stripped.startswith("  - "):
            text = stripped[4:].strip()
            bullet_level = 1
        else:
            text = stripped
            bullet_level = 0

        if not first_written:
            p = text_frame.paragraphs[0]
            first_written = True
        else:
            p = text_frame.add_paragraph()

        p.text = text
        p.level = bullet_level


def build_presentation(input_path: Path, output_path: Path):
    from pptx import Presentation

    md_text = input_path.read_text(encoding="utf-8")
    sections = parse_markdown_sections(md_text)

    if not sections:
        raise ValueError("No slide sections found. Ensure markdown uses '## Slide N — Title'.")

    prs = Presentation()

    add_title_slide(
        prs,
        title="Pantry Planner — Final Project Presentation",
        subtitle="Generated from markdown source",
    )

    for title, lines in sections:
        add_content_slide(prs, title=title, lines=lines)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    args = parser.parse_args()

    build_presentation(args.input, args.output)
    print(f"Created: {args.output}")


if __name__ == "__main__":
    main()
